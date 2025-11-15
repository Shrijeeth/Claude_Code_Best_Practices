import re
from pathlib import Path

import docx
import openpyxl
import pypdf
from pptx import Presentation


class DocumentProcessor:
    """Process various document types and extract text"""

    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def process_document(self, file_path: str) -> list[str]:
        """Process document and return text chunks"""
        file_path = Path(file_path)
        extension = file_path.suffix.lower()

        # Extract text based on file type
        if extension == ".pdf":
            text = self._extract_from_pdf(file_path)
        elif extension in [".docx", ".doc"]:
            text = self._extract_from_docx(file_path)
        elif extension in [".xlsx", ".xls"]:
            text = self._extract_from_excel(file_path)
        elif extension in [".pptx", ".ppt"]:
            text = self._extract_from_pptx(file_path)
        elif extension in [".txt", ".md"]:
            text = self._extract_from_text(file_path)
        else:
            raise ValueError(f"Unsupported file type: {extension}")

        # Split into chunks
        chunks = self._create_chunks(text)
        return chunks

    def _extract_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF"""
        text = ""
        with open(file_path, "rb") as file:
            pdf_reader = pypdf.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text

    def _extract_from_docx(self, file_path: Path) -> str:
        """Extract text from DOCX"""
        doc = docx.Document(file_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text

    def _extract_from_excel(self, file_path: Path) -> str:
        """Extract text from Excel"""
        workbook = openpyxl.load_workbook(file_path)
        text = ""
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                text += " ".join([str(cell) for cell in row if cell is not None]) + "\n"
        return text

    def _extract_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint"""
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def _extract_from_text(self, file_path: Path) -> str:
        """Extract text from plain text files"""
        with open(file_path, encoding="utf-8") as file:
            return file.read()

    def _create_chunks(self, text: str) -> list[str]:
        """Split text into overlapping chunks"""
        # Clean text
        text = re.sub(r"\s+", " ", text).strip()

        chunks = []
        start = 0
        text_length = len(text)

        while start < text_length:
            end = start + self.chunk_size

            # Try to break at sentence boundary
            if end < text_length:
                # Look for sentence ending
                chunk_text = text[start:end]
                last_period = max(
                    chunk_text.rfind(". "), chunk_text.rfind("? "), chunk_text.rfind("! ")
                )
                if last_period > self.chunk_size * 0.5:  # At least 50% of chunk size
                    end = start + last_period + 1

            chunks.append(text[start:end].strip())
            start = end - self.chunk_overlap

        return [chunk for chunk in chunks if chunk]
